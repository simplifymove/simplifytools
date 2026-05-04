# -*- coding: utf-8 -*-
"""
PDF Convert Engine
Handles format conversions: PDF to images, images to PDF, document conversions
"""

from typing import Dict, Any, List
from pathlib import Path
import PyPDF2
import fitz  # PyMuPDF
from PIL import Image
import zipfile
import os
import time
import io

# Enable HEIC support for Pillow
try:
    import pillow_heif
    pillow_heif.register_heif_opener()
except ImportError:
    pass  # pillow-heif not installed, HEIC support unavailable


class PdfConvertEngine:
    """PDF format conversion operations"""
    
    @staticmethod
    def pdf_to_image(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to JPG/PNG/TIFF images at specified DPI"""
        try:
            print(f"[PDF_TO_IMAGE] ===== START PDF TO IMAGE CONVERSION =====")
            print(f"[PDF_TO_IMAGE] Input paths: {input_paths}")
            print(f"[PDF_TO_IMAGE] Output path: {output_path}")
            print(f"[PDF_TO_IMAGE] Output path type: {type(output_path)}")
            print(f"[PDF_TO_IMAGE] Options: {options}")
            
            pdf_path = input_paths[0]
            output_format = options.get('format', 'jpg')  # jpg, png, tiff
            dpi = int(options.get('dpi', 150))
            page_mode = options.get('pageMode', 'all')  # all, selected
            page_range = options.get('pageRange', '')
            
            print(f"[PDF_TO_IMAGE] Opening PDF: {pdf_path}")
            doc = fitz.open(pdf_path)
            print(f"[PDF_TO_IMAGE] PDF opened successfully. Total pages: {len(doc)}")
            
            zoom = dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            
            pages_to_convert = []
            if page_mode == 'all':
                pages_to_convert = list(range(len(doc)))
            else:
                for part in page_range.split(','):
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages_to_convert.extend(range(start-1, end))
                    else:
                        pages_to_convert.append(int(part)-1)
            
            print(f"[PDF_TO_IMAGE] Pages to convert: {pages_to_convert}")
            
            output_dir = Path(output_path).parent
            print(f"[PDF_TO_IMAGE] Output directory: {output_dir}")
            print(f"[PDF_TO_IMAGE] Output directory exists: {output_dir.exists()}")
            
            output_files = []
            
            # Determine output format
            output_format_lower = output_format.lower()
            if output_format_lower == 'jpg':
                output_format_lower = 'jpeg'
                ext = 'jpg'
                needs_alpha = False
            elif output_format_lower == 'tiff':
                ext = 'tiff'
                needs_alpha = False
            elif output_format_lower == 'png':
                ext = 'png'
                needs_alpha = True  # PNG needs alpha channel
            else:
                ext = output_format_lower
                needs_alpha = False
            
            print(f"[PDF_TO_IMAGE] Output format: {output_format_lower}, Extension: {ext}")
            
            for page_num in pages_to_convert:
                page = doc[page_num]
                pix = page.get_pixmap(matrix=mat, alpha=needs_alpha)
                
                out_file = output_dir / f"page_{page_num+1}.{ext}"
                print(f"[PDF_TO_IMAGE] Creating page {page_num+1}: {out_file}")
                
                if output_format_lower == 'jpeg':
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(out_file, "JPEG", quality=95)
                elif output_format_lower == 'tiff':
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(out_file, "TIFF", compression='lzw')
                elif output_format_lower == 'png':
                    # PNG with alpha channel - use PIL with RGBA mode
                    img = Image.frombytes("RGBA", [pix.width, pix.height], pix.samples)
                    img.save(out_file, "PNG")
                else:
                    # Fallback for other formats
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(out_file, str(ext).upper())
                
                # Verify file was created
                if Path(out_file).exists():
                    size = Path(out_file).stat().st_size
                    print(f"[PDF_TO_IMAGE] [OK] File created: {out_file} ({size} bytes)")
                    output_files.append(str(out_file))
                else:
                    print(f"[PDF_TO_IMAGE] [FAIL] ERROR: File NOT created: {out_file}")
            
            doc.close()
            print(f"[PDF_TO_IMAGE] Closed PDF. Total files created: {len(output_files)}")
            
            # If single page return it, else zip
            if len(output_files) == 1:
                print(f"[PDF_TO_IMAGE] Single page detected, returning image directly")
                return output_files[0]
            else:
                try:
                    print(f"[PDF_TO_IMAGE] Multiple pages detected, creating ZIP from {len(output_files)} files")
                    print(f"[PDF_TO_IMAGE] Creating ZIP file: {output_path}")
                    
                    # Verify all files exist before zipping
                    for img_file in output_files:
                        if Path(img_file).exists():
                            size = Path(img_file).stat().st_size
                            print(f"[PDF_TO_IMAGE] File exists: {img_file} ({size} bytes)")
                        else:
                            print(f"[PDF_TO_IMAGE] ERROR: File not found: {img_file}")
                    
                    # Create ZIP in memory first for reliability
                    zip_buffer = io.BytesIO()
                    files_added = 0
                    
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_STORED) as zf:
                        for img_file in output_files:
                            img_path = Path(img_file)
                            if img_path.exists():
                                try:
                                    with open(img_file, 'rb') as f:
                                        file_data = f.read()
                                    filename = img_path.name
                                    zf.writestr(filename, file_data, compress_type=zipfile.ZIP_STORED)
                                    files_added += 1
                                    print(f"[PDF_TO_IMAGE] Added to ZIP: {filename} ({len(file_data)} bytes)")
                                except Exception as e:
                                    print(f"[PDF_TO_IMAGE] Failed to add {img_file}: {str(e)}")
                            else:
                                print(f"[PDF_TO_IMAGE] Skipping missing file: {img_file}")
                    
                    print(f"[PDF_TO_IMAGE] Files added to ZIP: {files_added}")
                    
                    # Write in-memory ZIP to disk
                    zip_buffer.seek(0)
                    zip_data = zip_buffer.getvalue()
                    print(f"[PDF_TO_IMAGE] ZIP buffer size: {len(zip_data)} bytes")
                    
                    with open(output_path, 'wb') as f:
                        bytes_written = f.write(zip_data)
                    print(f"[PDF_TO_IMAGE] Bytes written to disk: {bytes_written}")
                    
                    # Verify ZIP was created and is valid
                    if not Path(output_path).exists():
                        raise Exception("ZIP file was not created")
                    
                    zip_size = Path(output_path).stat().st_size
                    print(f"[PDF_TO_IMAGE] ZIP file size on disk: {zip_size} bytes")
                    
                    if zip_size == 0:
                        raise Exception("ZIP file is empty - no files were added")
                    
                    try:
                        with zipfile.ZipFile(output_path, 'r') as verify_zf:
                            file_list = verify_zf.namelist()
                            print(f"[PDF_TO_IMAGE] ZIP contains {len(file_list)} files: {file_list}")
                            test_result = verify_zf.testzip()
                            if test_result is not None:
                                raise Exception(f"ZIP file corrupt at: {test_result}")
                    except Exception as e:
                        raise Exception(f"ZIP file validation failed: {str(e)}")
                    
                    # Clean up temporary image files
                    for img_file in output_files:
                        try:
                            os.remove(img_file)
                        except Exception as e:
                            print(f"[WARNING] Failed to clean up {img_file}: {str(e)}")
                    
                    return output_path
                except Exception as e:
                    # Clean up on error
                    for img_file in output_files:
                        try:
                            os.remove(img_file)
                        except:
                            pass
                    raise Exception(f"Failed to create ZIP file: {str(e)}")
        except Exception as e:
            raise Exception(f"Failed to convert PDF to image: {str(e)}")
    
    @staticmethod
    def image_to_pdf(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert images to PDF with proper scaling and orientation"""
        try:
            # Use higher DPI for better quality (300 DPI instead of 72 DPI)
            # This provides 4x more detail when scaling
            dpi_factor = 300 / 72  # 4.167x scale factor
            
            # Standard PDF page size (8.5 x 11 inches at 300 DPI = 2550 x 3300 points)
            page_width = 612 * dpi_factor
            page_height = 792 * dpi_factor
            
            doc = fitz.open()
            
            for idx, img_path in enumerate(input_paths):
                try:
                    file_ext = Path(img_path).suffix.lower()
                    img = None
                    
                    # Validate file exists and is readable
                    if not Path(img_path).exists():
                        raise Exception(f"File not found: {img_path}")
                    
                    if not os.access(img_path, os.R_OK):
                        raise Exception(f"File is not readable: {img_path}")
                    
                    # Try to open image with PIL (works for most raster formats)
                    try:
                        img = Image.open(img_path)
                    except (OSError, IOError) as open_err:
                        # If PIL can't open it, try PyMuPDF for PDF/PostScript documents
                        if file_ext in ['.eps', '.ps', '.pdf']:
                            try:
                                # PyMuPDF can handle PDF and PostScript formats
                                # Use is_pdf=False to allow PostScript detection
                                doc_page = fitz.open(img_path)
                                if doc_page is None or len(doc_page) == 0:
                                    raise Exception(f"File appears to be empty or invalid: {img_path}")
                                    
                                page = doc_page[0]
                                # Render at higher zoom for EPS files (PostScript needs quality rendering)
                                if file_ext in ['.eps', '.ps']:
                                    pix = page.get_pixmap(matrix=fitz.Matrix(3, 3), alpha=False)
                                else:
                                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                                    
                                if pix is None:
                                    raise Exception(f"Could not render {file_ext.upper()} page to image")
                                    
                                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                                doc_page.close()
                                
                            except Exception as fitz_err:
                                error_msg = str(fitz_err)
                                if "Ghostscript" in error_msg or "not found" in error_msg.lower():
                                    raise Exception(f"Cannot convert {file_ext} file. PostScript/EPS files require proper Ghostscript support. Error: {error_msg}")
                                else:
                                    raise Exception(f"Failed to convert {file_ext} file: {error_msg}")
                        else:
                            raise Exception(f"Cannot open image file ({file_ext}). Supported formats: JPG, PNG, GIF, WebP, TIFF, HEIC, PDF, EPS. Error: {str(open_err)}")
                    
                    # Convert RGBA to RGB if needed (for JPEG compatibility)
                    if img.mode in ('RGBA', 'LA', 'P'):
                        bg = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        bg.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                        img = bg
                    elif img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    # Get image dimensions
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    page_aspect = page_width / page_height
                    
                    # Scale image to fit page while maintaining aspect ratio with margin
                    margin = 20 * dpi_factor
                    if aspect_ratio > page_aspect:
                        # Image is wider
                        new_width = page_width - margin
                        new_height = int(new_width / aspect_ratio)
                    else:
                        # Image is taller
                        new_height = page_height - margin
                        new_width = int(new_height * aspect_ratio)
                    
                    # Only resize if image is larger than target
                    if img.size[0] > new_width or img.size[1] > new_height:
                        img = img.resize((int(new_width), int(new_height)), Image.Resampling.LANCZOS)
                    
                    # Create page with high-DPI dimensions
                    page = doc.new_page(width=page_width, height=page_height)
                    
                    # Center the image on the page
                    x = (page_width - new_width) / 2
                    y = (page_height - new_height) / 2
                    
                    # Save image temporarily with high quality
                    import tempfile
                    temp_dir = tempfile.gettempdir()
                    temp_img_path = os.path.join(temp_dir, f'temp_img_{idx}_{int(time.time() * 1000)}.jpg')
                    # Save as JPEG with high quality to preserve details
                    img.save(temp_img_path, 'JPEG', quality=95)
                    
                    # Insert image at full dimensions for maximum quality
                    rect = fitz.Rect(x, y, x + new_width, y + new_height)
                    page.insert_image(rect, filename=temp_img_path)
                    
                    # Clean up temp file
                    try:
                        os.remove(temp_img_path)
                    except:
                        pass
                        
                except Exception as img_err:
                    raise Exception(f"Failed to process image {idx + 1} ({Path(img_path).name}): {str(img_err)}")
            
            if len(doc) == 0:
                raise Exception("No images were successfully converted")
            
            # Save PDF with high quality compression
            doc.save(output_path, deflate=True, garbage=4)
            doc.close()
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert images to PDF: {str(e)}")
    
    @staticmethod
    def pdf_to_text(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Extract text from PDF"""
        try:
            pdf_path = input_paths[0]
            page_range = options.get('pageRange', 'all')
            
            doc = fitz.open(pdf_path)
            text_content = []
            
            pages_to_extract = []
            if page_range == 'all':
                pages_to_extract = list(range(len(doc)))
            else:
                for part in page_range.split(','):
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages_to_extract.extend(range(start-1, end))
                    else:
                        pages_to_extract.append(int(part)-1)
            
            for page_num in pages_to_extract:
                page = doc[page_num]
                text = page.get_text()
                text_content.append(f"--- Page {page_num+1} ---\n{text}\n")
            
            doc.close()
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.writelines(text_content)
            
            return output_path
        except Exception as e:
            raise Exception(f"Failed to extract text: {str(e)}")
    
    @staticmethod
    def pdf_to_docx(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to Word DOCX format (with OCR support for scanned PDFs)"""
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            
            pdf_path = input_paths[0]
            page_range = options.get('pageRange', 'all')
            language = options.get('language', 'en')
            
            print(f"[PDF_TO_DOCX] Starting conversion: {pdf_path}", flush=True)
            
            doc_obj = fitz.open(pdf_path)
            document = Document()
            
            # Initialize OCR reader for fallback
            ocr_reader = None
            
            pages_to_extract = []
            if page_range == 'all':
                pages_to_extract = list(range(len(doc_obj)))
            else:
                for part in page_range.split(','):
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages_to_extract.extend(range(start-1, end))
                    else:
                        pages_to_extract.append(int(part)-1)
            
            for page_num in pages_to_extract:
                page = doc_obj[page_num]
                text = page.get_text()
                
                # Add page heading
                document.add_heading(f'Page {page_num + 1}', level=2)
                
                # Add content
                if text.strip():
                    print(f"[PDF_TO_DOCX] Text found on page {page_num + 1}", flush=True)
                    for paragraph_text in text.split('\n'):
                        if paragraph_text.strip():
                            p = document.add_paragraph(paragraph_text)
                            p.paragraph_format.space_before = Pt(6)
                            p.paragraph_format.space_after = Pt(6)
                else:
                    # Try OCR
                    print(f"[PDF_TO_DOCX] No text found, attempting OCR on page {page_num + 1}", flush=True)
                    try:
                        if ocr_reader is None:
                            import easyocr
                            print(f"[PDF_TO_DOCX] Initializing OCR reader", flush=True)
                            ocr_reader = easyocr.Reader([language], gpu=False, verbose=False)
                        
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                        import tempfile
                        temp_img = tempfile.NamedTemporaryFile(suffix='.png', delete=False).name
                        pix.save_png(temp_img)
                        
                        results = ocr_reader.readtext(temp_img)
                        ocr_text = '\n'.join([line[1] for line in results if line[1].strip()])
                        
                        if ocr_text.strip():
                            print(f"[PDF_TO_DOCX] OCR successful", flush=True)
                            for paragraph_text in ocr_text.split('\n'):
                                if paragraph_text.strip():
                                    p = document.add_paragraph(paragraph_text)
                                    p.paragraph_format.space_before = Pt(6)
                                    p.paragraph_format.space_after = Pt(6)
                        else:
                            document.add_paragraph('(No text detected)')
                        
                        try:
                            import os as os_module
                            os_module.unlink(temp_img)
                        except:
                            pass
                    except Exception as ocr_err:
                        print(f"[PDF_TO_DOCX] OCR failed: {str(ocr_err)}", flush=True)
                        document.add_paragraph('(Page with no extractable text)')
                
                document.add_paragraph()  # Add spacing between pages
            
            doc_obj.close()
            print(f"[PDF_TO_DOCX] Writing DOCX file", flush=True)
            document.save(output_path)
            return output_path
        except Exception as e:
            import traceback
            error_msg = f"Failed to convert PDF to DOCX: {str(e)}\n{traceback.format_exc()}"
            print(f"[PDF_TO_DOCX] ERROR: {error_msg}", flush=True)
            raise Exception(error_msg)
    
    @staticmethod
    def pdf_to_pptx(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to PowerPoint PPTX format"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            pdf_path = input_paths[0]
            page_range = options.get('pageRange', 'all')
            
            doc = fitz.open(pdf_path)
            prs = Presentation()
            
            # Set slide size
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            pages_to_extract = []
            if page_range == 'all':
                pages_to_extract = list(range(len(doc)))
            else:
                for part in page_range.split(','):
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages_to_extract.extend(range(start-1, end))
                    else:
                        pages_to_extract.append(int(part)-1)
            
            for page_num in pages_to_extract:
                page = doc[page_num]
                
                # Create blank slide
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Convert PDF page to image and add to slide
                mat = fitz.Matrix(1, 1)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                import tempfile
                temp_dir = tempfile.gettempdir()
                temp_img = os.path.join(temp_dir, f'slide_{page_num}_{int(time.time() * 1000)}.png')
                
                # Save pixmap to image file - use appropriate method for PyMuPDF
                try:
                    pix.save(temp_img)  # Newer PyMuPDF versions use save()
                except AttributeError:
                    pix.save_png(temp_img)  # Fallback for older versions
                
                # Add image to slide (scaled to fit)
                left = Inches(0.5)
                top = Inches(0.5)
                height = Inches(6.5)
                pic = slide.shapes.add_picture(temp_img, left, top, height=height)
                
                # Add page number
                txBox = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
                tf = txBox.text_frame
                tf.text = f"Page {page_num + 1}"
                p = tf.paragraphs[0]
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.RIGHT
                
                # Clean up
                try:
                    os.remove(temp_img)
                except:
                    pass
            
            doc.close()
            prs.save(output_path)
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert PDF to PPTX: {str(e)}")
    
    @staticmethod
    def pdf_to_xlsx(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to Excel XLSX format (extracts text and tables)"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            import pdfplumber
            
            pdf_path = input_paths[0]
            page_range = options.get('pageRange', 'all')
            
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet
            
            with pdfplumber.open(pdf_path) as pdf:
                pages_to_extract = []
                if page_range == 'all':
                    pages_to_extract = list(range(len(pdf.pages)))
                else:
                    for part in page_range.split(','):
                        if '-' in part:
                            start, end = map(int, part.split('-'))
                            pages_to_extract.extend(range(start-1, end))
                        else:
                            pages_to_extract.append(int(part)-1)
                
                for page_num in pages_to_extract:
                    page = pdf.pages[page_num]
                    
                    # Create sheet for this page
                    ws = wb.create_sheet(title=f"Page {page_num + 1}")
                    
                    # Extract tables if any
                    tables = page.extract_tables()
                    if tables:
                        row = 1
                        for table_idx, table in enumerate(tables):
                            if table_idx > 0:
                                row += 2  # Spacing between tables
                            
                            for table_row_idx, table_row in enumerate(table):
                                for col_idx, cell_val in enumerate(table_row):
                                    cell = ws.cell(row=row, column=col_idx + 1)
                                    cell.value = cell_val
                                    
                                    # Format header row
                                    if table_row_idx == 0:
                                        cell.font = Font(bold=True, color="FFFFFF")
                                        cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                                    
                                    cell.alignment = Alignment(wrap_text=True, vertical="top")
                                
                                row += 1
                    else:
                        # If no tables, extract text
                        text = page.extract_text()
                        if text:
                            lines = text.split('\n')
                            for line_idx, line in enumerate(lines):
                                cell = ws.cell(row=line_idx + 1, column=1)
                                cell.value = line
                                cell.alignment = Alignment(wrap_text=True)
            
            # Auto-adjust column widths
            for ws in wb.sheetnames:
                sheet = wb[ws]
                for column in sheet.columns:
                    max_length = 0
                    column_letter = column[0].column_letter
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    sheet.column_dimensions[column_letter].width = adjusted_width
            
            wb.save(output_path)
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert PDF to XLSX: {str(e)}")
    
    @staticmethod
    def pdf_to_html(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to HTML format (with OCR support for scanned PDFs)"""
        try:
            pdf_path = input_paths[0]
            page_range = options.get('pageRange', 'all')
            language = options.get('language', 'en')
            
            print(f"[PDF_TO_HTML] Starting conversion: {pdf_path}", flush=True)
            
            doc = fitz.open(pdf_path)
            
            # Initialize OCR reader for fallback
            ocr_reader = None
            
            pages_to_extract = []
            if page_range == 'all':
                pages_to_extract = list(range(len(doc)))
            else:
                for part in page_range.split(','):
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages_to_extract.extend(range(start-1, end))
                    else:
                        pages_to_extract.append(int(part)-1)
            
            html_content = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PDF to HTML Conversion</title>
    <style>
        body { font-family: Arial, sans-serif; line-height: 1.6; margin: 20px; }
        .page { page-break-after: always; margin-bottom: 40px; border: 1px solid #ddd; padding: 20px; }
        .page-number { text-align: right; color: #999; margin-top: 20px; }
        h1 { color: #333; }
        p { color: #666; }
    </style>
</head>
<body>
"""
            
            for page_num in pages_to_extract:
                page = doc[page_num]
                text = page.get_text()
                
                html_content += f'<div class="page">\n'
                html_content += f'<h1>Page {page_num + 1}</h1>\n'
                
                if text.strip():
                    print(f"[PDF_TO_HTML] Text found on page {page_num + 1}", flush=True)
                    for line in text.split('\n'):
                        if line.strip():
                            html_content += f'<p>{line.strip()}</p>\n'
                else:
                    # Try OCR
                    print(f"[PDF_TO_HTML] No text found, attempting OCR on page {page_num + 1}", flush=True)
                    try:
                        if ocr_reader is None:
                            import easyocr
                            print(f"[PDF_TO_HTML] Initializing OCR reader", flush=True)
                            ocr_reader = easyocr.Reader([language], gpu=False, verbose=False)
                        
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                        import tempfile
                        temp_img = tempfile.NamedTemporaryFile(suffix='.png', delete=False).name
                        pix.save_png(temp_img)
                        
                        results = ocr_reader.readtext(temp_img)
                        ocr_text = '\n'.join([line[1] for line in results if line[1].strip()])
                        
                        if ocr_text.strip():
                            print(f"[PDF_TO_HTML] OCR successful", flush=True)
                            for line in ocr_text.split('\n'):
                                if line.strip():
                                    html_content += f'<p>{line.strip()}</p>\n'
                        else:
                            html_content += '<p><em>(No text detected)</em></p>\n'
                        
                        try:
                            import os as os_module
                            os_module.unlink(temp_img)
                        except:
                            pass
                    except Exception as ocr_err:
                        print(f"[PDF_TO_HTML] OCR failed: {str(ocr_err)}", flush=True)
                        html_content += '<p><em>(Page with no extractable text)</em></p>\n'
                
                html_content += f'<div class="page-number">Page {page_num + 1}</div>\n'
                html_content += '</div>\n'
            
            html_content += """</body>
</html>"""
            
            doc.close()
            
            print(f"[PDF_TO_HTML] Writing HTML file", flush=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            return output_path
        except Exception as e:
            import traceback
            error_msg = f"Failed to convert PDF to HTML: {str(e)}\n{traceback.format_exc()}"
            print(f"[PDF_TO_HTML] ERROR: {error_msg}", flush=True)
            raise Exception(error_msg)
    
    @staticmethod
    def pdf_to_rtf(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to RTF format (with OCR support for scanned PDFs)"""
        try:
            pdf_path = input_paths[0]
            page_range = options.get('pageRange', 'all')
            language = options.get('language', 'en')
            
            print(f"[PDF_TO_RTF] Starting conversion: {pdf_path}", flush=True)
            print(f"[PDF_TO_RTF] Language for OCR: {language}", flush=True)
            
            doc = fitz.open(pdf_path)
            
            pages_to_extract = []
            if page_range == 'all':
                pages_to_extract = list(range(len(doc)))
            else:
                for part in page_range.split(','):
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages_to_extract.extend(range(start-1, end))
                    else:
                        pages_to_extract.append(int(part)-1)
            
            # Initialize OCR reader for fallback (lazy load on first use)
            ocr_reader = None
            
            rtf_content = r"{\rtf1\ansi\ansicpg1252\deff0{\fonttbl{\f0 Arial;}}" + "\n"
            rtf_content += r"{\colortbl;\red0\green0\blue0;}" + "\n"
            rtf_content += r"\viewkind4\uc1\pard\f0\fs20 " + "\n"
            
            for page_num in pages_to_extract:
                print(f"[PDF_TO_RTF] Processing page {page_num + 1}/{len(pages_to_extract)}", flush=True)
                page = doc[page_num]
                text = page.get_text()
                
                # RTF header for page
                rtf_content += r"\b PDF Page " + str(page_num + 1) + r"\b0 \par \par " + "\n"
                
                if text.strip():
                    print(f"[PDF_TO_RTF] [OK] Text found on page {page_num + 1}", flush=True)
                    for line in text.split('\n'):
                        if line.strip():
                            # Escape special RTF characters
                            line_escaped = line.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
                            rtf_content += line_escaped + r" \par " + "\n"
                else:
                    # No text found - try OCR
                    print(f"[PDF_TO_RTF] [NO_TEXT] No text found, attempting OCR on page {page_num + 1}...", flush=True)
                    
                    ocr_attempted = False
                    try:
                        # Lazy load OCR reader on first use
                        if ocr_reader is None:
                            try:
                                import easyocr
                                print(f"[PDF_TO_RTF] Importing easyocr...", flush=True)
                                print(f"[PDF_TO_RTF] Loading EasyOCR reader for language: {language}", flush=True)
                                ocr_reader = easyocr.Reader([language], gpu=False, verbose=False)
                                print(f"[PDF_TO_RTF] [OK] OCR reader initialized", flush=True)
                            except ImportError as ie:
                                print(f"[PDF_TO_RTF] [FAIL] EasyOCR import failed: {str(ie)}", flush=True)
                                raise Exception(f"EasyOCR not installed: {str(ie)}")
                            except Exception as e:
                                print(f"[PDF_TO_RTF] [FAIL] Failed to initialize OCR: {str(e)}", flush=True)
                                raise
                        
                        ocr_attempted = True
                        print(f"[PDF_TO_RTF] Rendering page to image for OCR...", flush=True)
                        
                        try:
                            # Render page to image for OCR (2x zoom for better quality)
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                            print(f"[PDF_TO_RTF] Image size: {pix.width}x{pix.height}", flush=True)
                        except Exception as pix_err:
                            print(f"[PDF_TO_RTF] [FAIL] Failed to render page: {str(pix_err)}", flush=True)
                            raise
                        
                        # Save to temp image
                        import tempfile
                        try:
                            temp_img = tempfile.NamedTemporaryFile(suffix='.png', delete=False).name
                            pix.save_png(temp_img)
                            print(f"[PDF_TO_RTF] Saved image to: {temp_img}", flush=True)
                        except Exception as save_err:
                            print(f"[PDF_TO_RTF] [FAIL] Failed to save image: {str(save_err)}", flush=True)
                            raise
                        
                        print(f"[PDF_TO_RTF] Running readtext OCR...", flush=True)
                        try:
                            # Run OCR
                            results = ocr_reader.readtext(temp_img)
                            print(f"[PDF_TO_RTF] OCR readtext completed, found {len(results)} regions", flush=True)
                        except Exception as ocr_read_err:
                            print(f"[PDF_TO_RTF] [FAIL] readtext failed: {str(ocr_read_err)}", flush=True)
                            raise
                        
                        # Extract text from OCR results
                        ocr_text = '\n'.join([line[1] for line in results if line[1].strip()])
                        print(f"[PDF_TO_RTF] Extracted {len(ocr_text)} characters from OCR", flush=True)
                        
                        if ocr_text.strip():
                            print(f"[PDF_TO_RTF] [OK] OCR successful", flush=True)
                            for line in ocr_text.split('\n'):
                                if line.strip():
                                    line_escaped = line.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
                                    rtf_content += line_escaped + r" \par " + "\n"
                        else:
                            print(f"[PDF_TO_RTF] [NO_TEXT] OCR returned no text", flush=True)
                            rtf_content += r"\i (No text detected by OCR) \i0 \par " + "\n"
                        
                        # Clean up temp file
                        try:
                            import os as os_module
                            os_module.unlink(temp_img)
                            print(f"[PDF_TO_RTF] Cleaned up temp file", flush=True)
                        except:
                            pass
                            
                    except Exception as ocr_err:
                        import traceback
                        tb = traceback.format_exc()
                        print(f"[PDF_TO_RTF] [FAIL] OCR exception: {str(ocr_err)}", flush=True)
                        print(f"[PDF_TO_RTF] Traceback: {tb}", flush=True)
                        rtf_content += r"\i (Page with no extractable text - OCR failed) \i0 \par " + "\n"
                
                rtf_content += r"\par \page " + "\n"
            
            rtf_content += "}"
            
            doc.close()
            
            print(f"[PDF_TO_RTF] Writing RTF file: {output_path}", flush=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(rtf_content)
            
            print(f"[PDF_TO_RTF] [OK] Conversion complete", flush=True)
            return output_path
            
        except Exception as e:
            import traceback
            error_msg = f"Failed to convert PDF to RTF: {str(e)}\n{traceback.format_exc()}"
            print(f"[PDF_TO_RTF] [ERROR] {error_msg}", flush=True)
            raise Exception(error_msg)
            
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert PDF to RTF: {str(e)}")
    
    @staticmethod
    def pdf_to_document(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to Word/PowerPoint/Excel/CSV (dispatcher)"""
        output_format = options.get('format', 'docx')  # docx, pptx, xlsx, html, rtf
        
        if output_format == 'docx':
            return PdfConvertEngine.pdf_to_docx(input_paths, output_path, options)
        elif output_format == 'pptx':
            return PdfConvertEngine.pdf_to_pptx(input_paths, output_path, options)
        elif output_format == 'xlsx':
            return PdfConvertEngine.pdf_to_xlsx(input_paths, output_path, options)
        elif output_format == 'html':
            return PdfConvertEngine.pdf_to_html(input_paths, output_path, options)
        elif output_format == 'rtf':
            return PdfConvertEngine.pdf_to_rtf(input_paths, output_path, options)
        else:
            raise ValueError(f"Format {output_format} not supported")
    
    @staticmethod
    def pdf_to_ebook(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert PDF to EPUB/MOBI/AZW3 (requires ebook-convert/calibre)"""
        try:
            pdf_path = input_paths[0]
            output_format = options.get('format', 'epub')
            
            # This requires ebook-convert (Calibre)
            import subprocess
            result = subprocess.run(
                ['ebook-convert', pdf_path, output_path],
                capture_output=True, text=True
            )
            
            if result.returncode != 0:
                raise Exception(f"ebook-convert failed: {result.stderr}")
            
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert PDF to ebook: {str(e)}")
    
    @staticmethod
    def ebook_to_pdf(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert EPUB/MOBI/AZW3 to PDF (requires ebook-convert/calibre)"""
        try:
            ebook_path = input_paths[0]
            
            # This requires ebook-convert (Calibre)
            import subprocess
            result = subprocess.run(
                ['ebook-convert', ebook_path, output_path],
                capture_output=True, text=True
            )
            
            if result.returncode != 0:
                raise Exception(f"ebook-convert failed: {result.stderr}")
            
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert ebook to PDF: {str(e)}")
    
    @staticmethod
    def document_to_pdf(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert Word/PowerPoint document to PDF (requires LibreOffice)"""
        try:
            import os
            import subprocess
            from pathlib import Path
            
            doc_path = input_paths[0]
            output_dir = str(Path(output_path).parent)
            output_filename = Path(output_path).name
            
            # LibreOffice converts files in-place with their original name but .pdf extension
            # So we need to get what LibreOffice will name it, then rename it to our desired output_path
            
            # Use LibreOffice headless with absolute path to ensure it works
            result = subprocess.run([
                'soffice', '--headless', '--convert-to', 'pdf', 
                '--outdir', output_dir,
                doc_path
            ], capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed (code {result.returncode}): {result.stderr}")
            
            # LibreOffice creates a file with the same name as input but .pdf extension
            input_name = Path(doc_path).stem  # Get filename without extension
            libreoffice_output = os.path.join(output_dir, f"{input_name}.pdf")
            
            # Check if LibreOffice created the file
            if not os.path.exists(libreoffice_output):
                raise Exception(f"LibreOffice did not create output file at {libreoffice_output}")
            
            # If our desired output path is different, rename it
            if libreoffice_output != output_path:
                if os.path.exists(output_path):
                    os.remove(output_path)
                os.rename(libreoffice_output, output_path)
            
            if not os.path.exists(output_path):
                raise Exception(f"Output file not created at {output_path}")
            
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert document: {str(e)}")
    
    @staticmethod
    def url_to_pdf(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert URL to PDF (requires Playwright/headless browser)"""
        try:
            url = options.get('url', '')
            if not url:
                raise ValueError("URL not provided")
            
            # Use Playwright for headless browser
            from playwright.sync_api import sync_playwright
            
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page()
                page.goto(url, wait_until='networkidle')
                page.pdf(path=output_path)
                browser.close()
            
            return output_path
        except Exception as e:
            raise Exception(f"Failed to convert URL to PDF: {str(e)}")
    
    @staticmethod
    def outlook_to_pdf(input_paths: List[str], output_path: str, options: Dict[str, Any]) -> str:
        """Convert Outlook MSG file to PDF"""
        try:
            msg_path = input_paths[0]
            
            # Try to use win32com on Windows for native support
            try:
                from win32com.client import Dispatch
                outlook = Dispatch("Outlook.Application")
                msg = outlook.CreateItemFromTemplate(msg_path)
                msg.SaveAs(output_path, 4)  # 4 = olSaveAsPDF
                return output_path
            except ImportError:
                # Fallback: extract text and create simple PDF
                raise Exception("outlook_to_pdf requires win32com on Windows")
        except Exception as e:
            raise Exception(f"Failed to convert Outlook to PDF: {str(e)}")
